from pdfminer.high_level import extract_text # pip install PDFMiner.six
from pptx import Presentation # pip install python-pptx
from docx import Document # pip install python-docx
from enum import Enum
import re

# Clean text of special characters and invisible characters
def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\x00-\x7F]+', '', text)
    return text.strip()

# Remove references (like [1], [2], etc.) from the text
def remove_references(text):
    # Regular expression for matching references like [1], [2], or (Smith, 2000)
    ref_pattern = r'(\[\d+\])|(\(\w+, \d{4}\))'
    cleaned_text = re.sub(ref_pattern, '', text)
    return cleaned_text

# Remove bibliography section
def remove_bibliography(text):
    # Look for the start of the bibliography or references section
    bibliography_keywords = ['references', 'bibliography']
    start_index = len(text)
    
    for keyword in bibliography_keywords:
        # Find the first occurrence of the keyword and remove everything after it
        keyword_position = text.lower().find(keyword)
        if keyword_position != -1:
            start_index = min(start_index, keyword_position)
    
    # Return the text before the bibliography section
    cleaned_text = text[:start_index]
    return cleaned_text

def extract_from_pdf(file):
    text = extract_text(file)
    cleaned_text = remove_references(remove_bibliography(clean_text(text)))
    return cleaned_text

def extract_from_docx(file):
    doc = Document(file)
    text = '\n'.join([para.text for para in doc.paragraphs])
    cleaned_text = remove_references(remove_bibliography(clean_text(text)))
    return cleaned_text

def extract_from_pptx(file):
    prs = Presentation(file)
    text = '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    cleaned_text = remove_references(remove_bibliography(clean_text(text)))
    return cleaned_text

class FileType(Enum):
    PDF = 1
    DOCX = 2
    PPTX = 3

def extract_from_file(file, filetype=FileType.DOCX):
    if filetype == FileType.PDF:
        return extract_from_pdf(file)
    elif filetype == FileType.DOCX:
        return extract_from_docx(file)
    elif filetype == FileType.PPTX:
        return extract_from_pptx(file)
    else:
        raise ValueError("Invalid filetype detected!")